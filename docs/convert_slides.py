#!/usr/bin/env python3
"""
Convert Slidev Markdown to editable PowerPoint.
Parses slides.md and creates a PPTX with actual text boxes (not images).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re


def parse_slidev_markdown(filepath):
    """Parse Slidev markdown into slides."""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators (--- with optional metadata)
    # Match --- followed by optional layout/class metadata
    raw_slides = re.split(r'\n---\n(?:layout:.*?\n(?:class:.*?\n)?---\n)?|\n---\n', content)
    
    slides = []
    for raw in raw_slides:
        if not raw.strip():
            continue
            
        # Skip frontmatter (first slide with theme, background, etc.)
        if 'theme:' in raw[:200] and 'highlighter:' in raw[:500]:
            # Extract title from frontmatter
            title_match = re.search(r'^title:\s*(.+)$', raw, re.MULTILINE)
            if title_match:
                slides.append({
                    'title': title_match.group(1).strip(),
                    'body': '',
                    'code_blocks': []
                })
            continue
        
        # Handle layout metadata at start of slide
        raw = re.sub(r'^layout:.*$', '', raw, flags=re.MULTILINE)
        raw = re.sub(r'^class:.*$', '', raw, flags=re.MULTILINE)
        
        # Extract code blocks first (preserve them)
        code_blocks = []
        code_pattern = r'```(\w*)\s*\n(.*?)```'
        for match in re.finditer(code_pattern, raw, re.DOTALL):
            lang = match.group(1) or 'text'
            code = match.group(2).strip()
            code_blocks.append({'lang': lang, 'code': code})
        
        # Replace code blocks with placeholders
        body_without_code = re.sub(code_pattern, '[CODE_BLOCK]', raw, flags=re.DOTALL)
        
        # Extract title (first # heading)
        title_match = re.search(r'^#\s+(.+)$', body_without_code, re.MULTILINE)
        title = title_match.group(1).strip() if title_match else ""
        
        # Clean up body
        body = body_without_code
        body = re.sub(r'^#\s+.+$', '', body, flags=re.MULTILINE)  # Remove title
        body = re.sub(r'^##\s+', '• ', body, flags=re.MULTILINE)  # Convert ## to bullets
        body = re.sub(r'^###\s+', '  ◦ ', body, flags=re.MULTILINE)  # Convert ### to sub-bullets
        
        # Remove Slidev/Vue-specific syntax
        body = re.sub(r'<v-clicks?>|</v-clicks?>', '', body)
        body = re.sub(r'<div[^>]*>', '', body)
        body = re.sub(r'</div>', '', body)
        body = re.sub(r'<span[^>]*>|</span>', '', body)
        body = re.sub(r'<a[^>]*>([^<]*)</a>', r'\1', body)  # Keep link text
        body = re.sub(r'<strong[^>]*>|</strong>', '', body)
        body = re.sub(r'<code[^>]*>|</code>', '', body)
        body = re.sub(r'<br\s*/?>', '\n', body)
        body = re.sub(r'<style>.*?</style>', '', body, flags=re.DOTALL)
        body = re.sub(r'v-click="?\d*"?', '', body)
        body = re.sub(r'@click="[^"]*"', '', body)
        body = re.sub(r'hover="[^"]*"', '', body)
        body = re.sub(r'\{[^}]*\}', '', body)  # Remove Vue bindings like {scale: 0.6}
        body = re.sub(r'class="[^"]*"', '', body)
        
        # Clean up markdown formatting
        body = re.sub(r'\*\*([^*]+)\*\*', r'\1', body)  # Bold
        body = re.sub(r'\*([^*]+)\*', r'\1', body)  # Italic
        body = re.sub(r'`([^`]+)`', r'\1', body)  # Inline code
        
        # Clean up extra whitespace
        body = re.sub(r'\n{3,}', '\n\n', body)
        body = body.strip()
        
        # Convert bullet points
        body = re.sub(r'^-\s+', '• ', body, flags=re.MULTILINE)
        
        if title or body or code_blocks:
            slides.append({
                'title': title,
                'body': body,
                'code_blocks': code_blocks
            })
    
    return slides


def create_pptx(slides, output_path):
    """Create a PPTX with editable text."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, slide_data in enumerate(slides):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        if slide_data['title']:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data['title']
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
        
        current_top = Inches(1.4)
        
        # Add body text
        if slide_data['body']:
            # Split body by code block placeholders
            parts = slide_data['body'].split('[CODE_BLOCK]')
            code_idx = 0
            
            for j, part in enumerate(parts):
                part = part.strip()
                if part:
                    # Calculate height based on content
                    lines = part.count('\n') + 1
                    height = min(Inches(0.3 * lines), Inches(2.5))
                    
                    body_box = slide.shapes.add_textbox(
                        Inches(0.5), current_top, Inches(12.3), height
                    )
                    body_frame = body_box.text_frame
                    body_frame.word_wrap = True
                    
                    # Add text line by line
                    lines_list = part.split('\n')
                    for k, line in enumerate(lines_list):
                        if k == 0:
                            para = body_frame.paragraphs[0]
                        else:
                            para = body_frame.add_paragraph()
                        para.text = line
                        para.font.size = Pt(14)
                        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    
                    current_top += height + Inches(0.1)
                
                # Add code block after this text part
                if code_idx < len(slide_data['code_blocks']):
                    code_data = slide_data['code_blocks'][code_idx]
                    code_text = code_data['code']
                    lang = code_data['lang']
                    
                    # Calculate height based on code lines
                    code_lines = code_text.count('\n') + 1
                    code_height = min(Inches(0.22 * code_lines + 0.3), Inches(4.5))
                    
                    # Add code block with background styling hint
                    code_box = slide.shapes.add_textbox(
                        Inches(0.5), current_top, Inches(12.3), code_height
                    )
                    code_frame = code_box.text_frame
                    code_frame.word_wrap = True
                    
                    # Add language label and code
                    header = f"[{lang}]" if lang and lang != 'text' else ""
                    full_code = f"{header}\n{code_text}" if header else code_text
                    
                    code_lines_list = full_code.split('\n')
                    for k, line in enumerate(code_lines_list):
                        if k == 0:
                            para = code_frame.paragraphs[0]
                        else:
                            para = code_frame.add_paragraph()
                        para.text = line
                        para.font.size = Pt(10)
                        para.font.name = "Courier New"
                        para.font.color.rgb = RGBColor(0x2d, 0x2d, 0x2d)
                    
                    current_top += code_height + Inches(0.15)
                    code_idx += 1
        
        # If no body but has code blocks, add them
        elif slide_data['code_blocks']:
            for code_data in slide_data['code_blocks']:
                code_text = code_data['code']
                lang = code_data['lang']
                
                code_lines = code_text.count('\n') + 1
                code_height = min(Inches(0.22 * code_lines + 0.3), Inches(5))
                
                code_box = slide.shapes.add_textbox(
                    Inches(0.5), current_top, Inches(12.3), code_height
                )
                code_frame = code_box.text_frame
                code_frame.word_wrap = True
                
                header = f"[{lang}]" if lang and lang != 'text' else ""
                full_code = f"{header}\n{code_text}" if header else code_text
                
                code_lines_list = full_code.split('\n')
                for k, line in enumerate(code_lines_list):
                    if k == 0:
                        para = code_frame.paragraphs[0]
                    else:
                        para = code_frame.add_paragraph()
                    para.text = line
                    para.font.size = Pt(10)
                    para.font.name = "Courier New"
                    para.font.color.rgb = RGBColor(0x2d, 0x2d, 0x2d)
                
                current_top += code_height + Inches(0.15)
    
    prs.save(output_path)
    print(f"✓ Created {output_path} with {len(slides)} editable slides")


if __name__ == "__main__":
    import sys
    
    input_file = sys.argv[1] if len(sys.argv) > 1 else "slides.md"
    output_file = sys.argv[2] if len(sys.argv) > 2 else "App_Presentation.pptx"
    
    print(f"Converting {input_file} to {output_file}...")
    slides = parse_slidev_markdown(input_file)
    print(f"Found {len(slides)} slides")
    create_pptx(slides, output_file)
